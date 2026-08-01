import os
import json
import datetime
from typing import List, Dict, Any, Optional
import numpy as np

# Optional imports with fallbacks to avoid application crashes
try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import chromadb
except ImportError:
    chromadb = None

try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None

try:
    from openai import OpenAI
except ImportError:
    OpenAI = None

from app.config import settings

# Initialize SentenceTransformer model for offline embeddings
embedding_model = None
if SentenceTransformer:
    try:
        # Load small, efficient embedding model (runs fast on CPU)
        embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    except Exception as e:
        print(f"Error loading SentenceTransformer: {e}")

# Initialize ChromaDB client
chroma_client = None
collection = None
if chromadb:
    try:
        os.makedirs(settings.CHROMA_DB_DIR, exist_ok=True)
        chroma_client = chromadb.PersistentClient(path=settings.CHROMA_DB_DIR)
        # Create or get standard collection
        collection = chroma_client.get_or_create_collection(name="genesis_brain")
    except Exception as e:
        print(f"Error initializing ChromaDB: {e}")

# In-memory vector store fallback in case Chroma DB fails
fallback_vector_store = []  # list of dicts: {"user_id": int, "file_id": int, "text": str, "vector": np.ndarray}

def get_embedding(text: str) -> List[float]:
    """Generates an embedding vector for the given text."""
    if settings.OPENAI_API_KEY and OpenAI:
        try:
            client = OpenAI(api_key=settings.OPENAI_API_KEY)
            response = client.embeddings.create(
                input=[text],
                model="text-embedding-ada-002"
            )
            return response.data[0].embedding
        except Exception as e:
            print(f"OpenAI embedding error: {e}, falling back to sentence-transformers")
    
    if embedding_model:
        vector = embedding_model.encode(text)
        return vector.tolist()
    
    # Absolute fallback (random vector for testing)
    return np.random.uniform(-1, 1, 384).tolist()

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Parses text content from files of different formats."""
    text = ""
    file_type = file_type.upper()
    
    if not os.path.exists(file_path):
        return "File not found."

    if file_type == "TXT":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
            
    elif file_type == "PDF":
        if PyPDF2:
            try:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                text = f"Error reading PDF: {e}"
        else:
            text = "PDF reader library is not available."
            
    elif file_type == "DOCX":
        if docx:
            try:
                doc = docx.Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                text = f"Error reading DOCX: {e}"
        else:
            text = "DOCX reader library is not available."
            
    elif file_type == "PPT" or file_type == "PPTX":
        if pptx:
            try:
                prs = pptx.Presentation(file_path)
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_texts.append(shape.text)
                text = "\n".join(slide_texts)
            except Exception as e:
                text = f"Error reading PPT: {e}"
        else:
            text = "PPT reader library is not available."
    else:
        text = "Unsupported file type."
        
    return text

def chunk_text(text: str, chunk_size: int = 600, overlap: int = 150) -> List[str]:
    """Splits large text blocks into overlapping chunks."""
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk_words = words[i:i + chunk_size]
        chunks.append(" ".join(chunk_words))
        i += chunk_size - overlap
        if len(chunk_words) < chunk_size:
            break
    return chunks

class RAGService:
    @staticmethod
    def process_and_index_file(user_id: int, file_id: int, file_path: str, file_type: str):
        """Extracts text, chunks it, computes embeddings, and stores in Vector DB."""
        raw_text = extract_text_from_file(file_path, file_type)
        if not raw_text or raw_text.startswith("Error") or raw_text.startswith("Unsupported"):
            print(f"Skipping indexing due to read errors: {raw_text[:100]}")
            return
            
        chunks = chunk_text(raw_text)
        print(f"Indexing {len(chunks)} chunks for user_id={user_id}, file_id={file_id}")
        
        for idx, chunk in enumerate(chunks):
            embedding = get_embedding(chunk)
            chunk_id = f"user_{user_id}_file_{file_id}_chunk_{idx}"
            metadata = {
                "user_id": user_id,
                "file_id": file_id,
                "filename": os.path.basename(file_path)
            }
            
            if collection:
                try:
                    collection.add(
                        ids=[chunk_id],
                        embeddings=[embedding],
                        metadatas=[metadata],
                        documents=[chunk]
                    )
                except Exception as e:
                    print(f"Chroma add error: {e}. Storing in fallback store.")
                    fallback_vector_store.append({
                        "user_id": user_id,
                        "file_id": file_id,
                        "text": chunk,
                        "vector": np.array(embedding)
                    })
            else:
                fallback_vector_store.append({
                    "user_id": user_id,
                    "file_id": file_id,
                    "text": chunk,
                    "vector": np.array(embedding)
                })

    @staticmethod
    def retrieve_context(user_id: int, file_id: Optional[int], query: str, top_k: int = 4) -> List[str]:
        """Retrieves relevant text fragments matching the user query."""
        query_vector = get_embedding(query)
        
        # Use Chroma if active
        if collection:
            try:
                where_filter = {"user_id": user_id}
                if file_id is not None:
                    where_filter["file_id"] = file_id
                    
                results = collection.query(
                    query_embeddings=[query_vector],
                    n_results=top_k,
                    where=where_filter
                )
                if results and 'documents' in results and results['documents']:
                    return results['documents'][0]
            except Exception as e:
                print(f"Chroma DB query failed: {e}. Falling back to in-memory matching.")
        
        # In-memory retrieval fallback
        matches = []
        for item in fallback_vector_store:
            if item["user_id"] == user_id:
                if file_id is None or item["file_id"] == file_id:
                    # Calculate Cosine Similarity
                    norm_a = np.linalg.norm(item["vector"])
                    norm_b = np.linalg.norm(query_vector)
                    sim = np.dot(item["vector"], query_vector) / (norm_a * norm_b) if norm_a > 0 and norm_b > 0 else 0
                    matches.append((sim, item["text"]))
                    
        matches.sort(key=lambda x: x[0], reverse=True)
        return [text for sim, text in matches[:top_k]]

    @staticmethod
    def query_document(user_id: int, file_id: Optional[int], query: str, mode: str = "standard") -> str:
        """Runs the RAG prompt pipeline with various summary / output modes."""
        contexts = RAGService.retrieve_context(user_id, file_id, query)
        context_text = "\n\n".join(contexts) if contexts else "No relevant document text was found."
        
        system_prompt = (
            "You are Genesis Brain, the RAG QA engine for Project Genesis.\n"
            "You answer questions based strictly on the retrieved document context below.\n"
            f"Mode selected: {mode}.\n"
            "Context:\n"
            f"{context_text}\n"
        )
        
        user_prompt = f"Question: {query}"
        
        # Customize output according to mode
        if mode == "flashcards":
            user_prompt += "\nFormat the response as a list of Question-and-Answer flashcards (Q: ... A: ...)."
        elif mode == "quiz":
            user_prompt += "\nGenerate a 5-question multiple choice quiz with answers key based on this content."
        elif mode == "beginner":
            user_prompt += "\nExplain this topic in extremely simple terms, as if explaining to a 10-year old."

        # Query LLM
        if settings.OPENAI_API_KEY and OpenAI:
            try:
                client = OpenAI(api_key=settings.OPENAI_API_KEY)
                completion = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ]
                )
                return completion.choices[0].message.content
            except Exception as e:
                print(f"OpenAI completion error: {e}")
                
        # LLM Simulation Fallback
        return RAGService._simulate_llm_response(query, context_text, mode)

    @staticmethod
    def _simulate_llm_response(query: str, context: str, mode: str) -> str:
        """Simulates RAG responses when OpenAI key is absent."""
        if context == "No relevant document text was found.":
            return "I couldn't find any relevant details in your uploaded documents. Please try uploading a document first."
            
        sentences = [s.strip() for s in context.split(".") if len(s.strip()) > 10]
        snippet = sentences[0] if sentences else context[:300]
        snippet2 = sentences[1] if len(sentences) > 1 else context[300:600]
        
        if mode == "flashcards":
            return (
                f"### 📇 Flashcards (AI Simulation)\n\n"
                f"**Card 1:**\n"
                f"Q: What is the main theme mentioned in this section?\n"
                f"A: The text focuses on: \"{snippet[:120]}...\"\n\n"
                f"**Card 2:**\n"
                f"Q: Can we extract details about the core concepts?\n"
                f"A: Yes, the source mentions: \"{snippet2[:120]}...\""
            )
        elif mode == "quiz":
            return (
                f"### 📝 Quick Quiz (AI Simulation)\n\n"
                f"**1. What primary topic does the text address?**\n"
                f"   a) Artificial Intel\n"
                f"   b) \"{snippet[:40]}\"\n"
                f"   c) Quantum Physics\n"
                f"   *Answer: b*\n\n"
                f"**2. According to details in the context, which of the following is true?**\n"
                f"   a) It mentions: \"{snippet2[:50]}\"\n"
                f"   b) It is completely false\n"
                f"   *Answer: a*"
            )
        elif mode == "beginner":
            return (
                f"### 👶 Explain Like I'm 5 (AI Simulation)\n\n"
                f"Imagine this topic is like playing with Lego blocks. The text is basically saying:\n\n"
                f"\"{snippet}.\"\n\n"
                f"Simply put, it means that parts work together, just like building a Lego castle!"
            )
        
        # Standard fallback response
        return (
            f"🤖 *Genesis Brain (Simulated Response)*\n\n"
            f"Based on the documents, here is the relevant information:\n\n"
            f"1. **Primary Insight**: \"{snippet}.\"\n"
            f"2. **Secondary Insight**: \"{snippet2}.\"\n\n"
            f"Let me know if you would like me to compile a quiz, explanation, or flashcards instead!"
        )


class StudyPlannerAI:
    @staticmethod
    def generate_plan(subjects: List[str], exam_date: datetime.datetime, available_hours: float) -> Dict[str, Any]:
        """Generates modular daily study schedule and revision milestones."""
        days_remaining = max(1, (exam_date.date() - datetime.datetime.utcnow().date()).days)
        
        system_prompt = (
            "You are a study planner AI. Generate a structured JSON study plan.\n"
            "Output must be a clean JSON object containing 'daily_schedule' (list of tasks/topics), "
            "'weekly_schedule' (list of key goals per week), and 'revision_milestones' (list of dates with review goals)."
        )
        user_prompt = (
            f"Subjects: {', '.join(subjects)}\n"
            f"Exam Date: {exam_date.strftime('%Y-%m-%d')} ({days_remaining} days away)\n"
            f"Daily hours: {available_hours} hours"
        )
        
        if settings.OPENAI_API_KEY and OpenAI:
            try:
                client = OpenAI(api_key=settings.OPENAI_API_KEY)
                completion = client.chat.completions.create(
                    model="gpt-4o-mini",
                    response_format={"type": "json_object"},
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ]
                )
                return json.loads(completion.choices[0].message.content)
            except Exception as e:
                print(f"Study planner LLM error: {e}")
                
        # Smart simulation fallback
        daily = []
        for idx, sub in enumerate(subjects):
            daily.append({
                "time": f"09:00 AM - 10:30 AM",
                "subject": sub,
                "activity": f"Read Chapter {idx+1} and take core summary notes"
            })
            daily.append({
                "time": f"02:00 PM - 03:00 PM",
                "subject": sub,
                "activity": f"Work on practice questions and review flashcards"
            })
            
        weekly = []
        weeks = max(1, days_remaining // 7)
        for w in range(1, weeks + 1):
            weekly.append({
                "week": f"Week {w}",
                "focus": f"Master fundamental concepts for {', '.join(subjects[:2])}",
                "milestone": f"Complete practice exam paper {w}"
            })
            
        revision = []
        midpoint = days_remaining // 2
        final_week = max(1, days_remaining - 3)
        
        revision.append({
            "days_left": midpoint,
            "title": "Midterm Checkpoint",
            "goal": "Re-evaluate and solve weak spots in subjects."
        })
        revision.append({
            "days_left": 3,
            "title": "Final Sprint Review",
            "goal": "Active recall and timed sample tests. No new material."
        })
        
        return {
            "daily_schedule": daily,
            "weekly_schedule": weekly,
            "revision_milestones": revision
        }


class ChatAssistantAI:
    @staticmethod
    def answer_query(user_profile: Dict[str, Any], query: str, context_notes: List[str], recent_tasks: List[str]) -> str:
        """Integrates user's tasks, profile, and notes to provide conversational advice."""
        profile_summary = f"User Interests: {user_profile.get('interests', [])}. Goal: {user_profile.get('daily_hours_goal')} hours study/productivity daily."
        tasks_text = ", ".join(recent_tasks) if recent_tasks else "No active tasks"
        notes_text = "\n- ".join(context_notes) if context_notes else "No recent notes"
        
        system_prompt = (
            "You are Genesis OS, a helpful, intelligent personal AI companion.\n"
            "You have access to the user's profile, recent notes, and active tasks.\n"
            "Help them plan, solve questions, structure study schedules, or summarize.\n"
            "Keep the tone futuristic, inspiring, and direct.\n\n"
            f"User Profile: {profile_summary}\n"
            f"Active Tasks: {tasks_text}\n"
            f"User Notes Context:\n{notes_text}\n"
        )
        
        if settings.OPENAI_API_KEY and OpenAI:
            try:
                client = OpenAI(api_key=settings.OPENAI_API_KEY)
                completion = client.chat.completions.create(
                    model="gpt-4o-mini",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": query}
                    ]
                )
                return completion.choices[0].message.content
            except Exception as e:
                print(f"Chat AI error: {e}")
                
        # Smart Simulation Fallback
        query_lower = query.lower()
        if "study" in query_lower:
            return (
                f"🌌 **Genesis OS Assistant**\n\n"
                f"Looking at your interests ({', '.join(user_profile.get('interests', []))}), I recommend prioritizing a structured block today.\n\n"
                f"You have some active tasks like *{tasks_text[:60]}*. I suggest setting a Pomodoro timer for 45 minutes to conquer these first."
            )
        elif "task" in query_lower or "important" in query_lower:
            return (
                f"🌌 **Genesis OS Assistant**\n\n"
                f"Your most critical tasks right now: *{tasks_text}*.\n\n"
                f"Action recommendation: Focus on completing the oldest task first to clear your backlog."
            )
        elif "note" in query_lower or "summarize" in query_lower:
            return (
                f"🌌 **Genesis OS Assistant**\n\n"
                f"Here is a quick summary of your recent notes:\n"
                f"*{notes_text}*\n\n"
                f"Let me know if you want me to expand on any topic!"
            )
            
        return (
            f"🌌 **Genesis OS Companion**\n\n"
            f"Greetings! I am scanning your digital landscape. Your goals indicate a target of {user_profile.get('daily_hours_goal')} productivity hours.\n\n"
            f"You asked: *\"{query}\"*\n\n"
            f"My response: I recommend aligning this with your interests ({', '.join(user_profile.get('interests', []))}). Let me know if we should add an event to your calendar or schedule a focus session!"
        )
